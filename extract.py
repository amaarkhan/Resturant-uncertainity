import os
import sys
import PyPDF2
from docx import Document
from pptx import Presentation

sys.stdout.reconfigure(encoding='utf-8')

CONTEXT_DIR = r"C:\Users\Lenovo\Documents\SEMESTER8\AIP\MVP\Previous Context"

def extract_pdf(path):
    text = ""
    try:
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text += page.extract_text() + "\n"
    except Exception as e:
        text = f"Error reading PDF: {e}"
    return text

def extract_docx(path):
    text = ""
    try:
        doc = Document(path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        text = f"Error reading DOCX: {e}"
    return text

def extract_pptx(path):
    text = ""
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        text = f"Error reading PPTX: {e}"
    return text

files = os.listdir(CONTEXT_DIR)
for file in files:
    path = os.path.join(CONTEXT_DIR, file)
    print(f"\n{'='*80}")
    print(f"FILE: {file}")
    print(f"{'='*80}\n")
    if file.endswith(".pdf"):
        print(extract_pdf(path))
    elif file.endswith(".docx"):
        print(extract_docx(path))
    elif file.endswith(".pptx"):
        print(extract_pptx(path))

